import os
from dotenv import load_dotenv
import streamlit as st

from langchain.schema import Document
from langchain.text_splitter import CharacterTextSplitter
from langchain.vectorstores import FAISS
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.document_loaders import PyPDFLoader

from docx import Document as DocxDocument
from pptx import Presentation

from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
import torch

# Load environment variables (if needed)
load_dotenv()

st.title("📄 Multi-Format Notes Chatbot (Flan-T5 Base - Local)")

uploaded_file = st.file_uploader("Upload your Notes File (PDF, DOCX, PPTX)", type=["pdf", "docx", "pptx"])

# ✅ Load the Flan-T5 Base model and tokenizer locally
@st.cache_resource
def load_model():
    model_name = "google/flan-t5-base"
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    model = AutoModelForSeq2SeqLM.from_pretrained(model_name)
    return tokenizer, model

tokenizer, model = load_model()

# ✅ Generate response using Flan-T5
def generate_answer(question, context):
    prompt = f"question: {question} context: {context}"
    inputs = tokenizer(prompt, return_tensors="pt", truncation=True, padding=True)
    outputs = model.generate(**inputs, max_length=512)
    return tokenizer.decode(outputs[0], skip_special_tokens=True)

# ✅ File reading functions
def extract_text_from_docx(path):
    doc = DocxDocument(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# ✅ Main logic
if uploaded_file is not None:
    file_extension = uploaded_file.name.split('.')[-1]
    file_path = f"temp.{file_extension}"

    with open(file_path, "wb") as f:
        f.write(uploaded_file.read())

    if file_path.endswith(".pdf"):
        loader = PyPDFLoader(file_path)
        documents = loader.load()
    elif file_path.endswith(".docx"):
        text = extract_text_from_docx(file_path)
        documents = [Document(page_content=text)]
    elif file_path.endswith(".pptx"):
        text = extract_text_from_pptx(file_path)
        documents = [Document(page_content=text)]
    else:
        st.error("Unsupported file format.")
        st.stop()

    text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    docs = text_splitter.split_documents(documents)

    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    db = FAISS.from_documents(docs, embeddings)

    question = st.text_input("Ask a question about the document:")

    if question:
        related_docs = db.similarity_search(question, k=3)
        combined_context = " ".join([doc.page_content for doc in related_docs])
        answer = generate_answer(question, combined_context)
        st.subheader("Answer:")
        st.write(answer)
